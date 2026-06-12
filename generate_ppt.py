from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Apply a dark theme approach using slide background color if possible
# But basic python-pptx templates are white. We will add a dark rectangle or simple text

slides_data = [
    {
        "title": "SKILLNOVA",
        "subtitle": "AI Driven Personalized Learning & Skill Gap Analyzer\n\nPresented By: [Your Name]"
    },
    {
        "title": "The Problem: The Industry Skill Gap",
        "content": "• Traditional learning is static; students don't know exactly what the industry demands today.\n• Existing platforms offer generic roadmaps without individually testing the user's actual capabilities."
    },
    {
        "title": "The Solution: SkillNova AI Ecosystem",
        "content": "• A deeply intelligent 'Career Twin' ecosystem.\n• It analyzes your exact resume/CGPA, maps your weaknesses, and tests you in real-time."
    },
    {
        "title": "Technical Architecture",
        "content": "• Frontend: React, Vite, Tailwind CSS, Framer Motion\n• Backend Auth: Node.js, Express, JWT Security\n• AI Engine: FastAPI (Python), Google Gemini API"
    },
    {
        "title": "User Execution Flow",
        "content": "Registration → Dashboard → Profile Sync → Skill Gap Analysis → Dynamic Roadmap → Practice Battle"
    },
    {
        "title": "Secure Access Node",
        "content": "• End-to-End hashed passwords using BCrypt.\n• Secure stateless sessions with JSON Web Tokens (JWT)."
    },
    {
        "title": "The Bridge Factor",
        "content": "• The central routing matrix of the application.\n• Users can launch AI Analysis, Roadmaps, or Battles directly from this seamless hub."
    },
    {
        "title": "Module 1 - Resume Parsing & Academic Matrix",
        "content": "• Upload a PDF or input academic CGPA/grades.\n• Extracted data is sent to the Python API to find exact technical keywords and frameworks."
    },
    {
        "title": "Module 2 - Generative AI Analysis",
        "content": "• Highly complex text structures are sent to Google Gemini 2.0.\n• The AI calculates an 'Industry Readiness Score' based on market trends and the user's specific skills."
    },
    {
        "title": "Fallback Logic Flow (Resilience Engine)",
        "content": "• If the Google AI goes down, we built a fallback 'Lexicon Hub' using Regex patterns.\n• The system never breaks during a live exam."
    },
    {
        "title": "Module 3 - Visual Skill Gap Detection",
        "content": "• Translates the AI JSON response into visual radar/bar charts.\n• Identifies exactly which frameworks the user needs to learn next."
    },
    {
        "title": "Module 4 - Dynamic Roadmap Generation",
        "content": "• Generates a non-linear pathway based securely on the detected skill gaps.\n• Provides recommended projects and milestones to reach 'S-Tier Readiness'."
    },
    {
        "title": "Module 5 - The Practice Battle Arena",
        "content": "• A gamified, timed testing environment fetching randomized generative AI questions.\n• Uses strict algorithmic matching to verify correct technical answers in seconds."
    },
    {
        "title": "Evaluation & Accuracy Scoring",
        "content": "• Converts Battle Accuracy into Core Metrics: Correctness, Clarity, and Depth.\n• Provides a highly accurate synthesis summary of your capabilities."
    },
    {
        "title": "Conclusion & Live Deployment",
        "content": "• SkillNova perfectly bridges the gap between static education and dynamic industry requirements.\n• Ready for live demonstration."
    }
]

# Title Slide Layout
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = slides_data[0]["title"]
subtitle.text = slides_data[0]["subtitle"]

# Bullet Slide Layout
bullet_slide_layout = prs.slide_layouts[1]

for slide_data in slides_data[1:]:
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = slide_data["title"]
    
    tf = body_shape.text_frame
    tf.text = slide_data["content"]
    
    # Add a placeholder box for screenshots
    left = Inches(5.5)
    top = Inches(2.0)
    width = Inches(4.0)
    height = Inches(4.5)
    
    # Adjust textbox size so they don't overlap with placeholder
    body_shape.width = Inches(4.5)
    
    # Add rectangle
    shape = shapes.add_shape(
        1, left, top, width, height  # 1 is MSO_SHAPE.RECTANGLE
    )
    shape.text = "[ PLACE SCREENSHOT HERE ]"
    
    # Optional: style the rectangle
    from pptx.dml.color import RGBColor
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(50, 50, 50)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].font.size = Pt(18)

prs.save('SkillNova_Presentation.pptx')
print("Presentation generated successfully: SkillNova_Presentation.pptx")
